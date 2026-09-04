import io
import zipfile

import openpyxl
from openpyxl.drawing.image import Image as XLImage
from PIL import Image
from pptx import Presentation

from extractors.structured_figures import (
    extract_docx_figures,
    extract_hwpx_figures,
    extract_pdf_figures,
    extract_pptx_figures,
    extract_xlsx_figures,
)


def _png() -> bytes:
    image = Image.new("RGB", (180, 100), "white")
    output = io.BytesIO()
    image.save(output, "PNG")
    return output.getvalue()


def _analyzer(_data, _mime, _filename, context):
    return {
        "kind": "table_image",
        "ocrText": "학문기초교양",
        "description": f"문맥을 반영한 필터 설명: {context[:30]}",
        "contextMatch": "supports",
        "confidence": 0.98,
    }


def _assert_contract(result, method):
    assert "[그림 1]" in result["text"]
    assert "[그림 설명]" in result["text"]
    assert result["figures"][0]["matchMethod"] == method
    assert result["figures"][0]["analysis"]["contextMatch"] == "supports"
    assert result["derived_assets"][0]["kind"] == "attachment_document_image"
    assert result["figure_contents"][0]["text"].startswith("[그림 1]\n")


def test_docx_figure_uses_paragraph_relationship_and_context(tmp_path):
    source = io.BytesIO()
    with zipfile.ZipFile(source, "w") as zf:
        zf.writestr("word/document.xml", """<w:document xmlns:w="urn:w" xmlns:a="urn:a" xmlns:r="urn:r"><w:body><w:p><w:r><w:t>학문기초교양 3학점 이수</w:t></w:r><a:blip r:embed="rId1"/></w:p></w:body></w:document>""")
        zf.writestr("word/_rels/document.xml.rels", """<Relationships><Relationship Id="rId1" Target="media/image1.png"/></Relationships>""")
        zf.writestr("word/media/image1.png", _png())

    result = extract_docx_figures(source.getvalue(), "", tmp_path, _analyzer)

    _assert_contract(result, "docx_paragraph_relationship")
    assert "학문기초교양" in result["figures"][0]["context"]


def test_hwpx_figure_uses_binary_item_reference(tmp_path):
    source = io.BytesIO()
    with zipfile.ZipFile(source, "w") as zf:
        zf.writestr("Contents/content.hpf", """<package><item id="image1" href="BinData/image1.png"/></package>""")
        zf.writestr("Contents/section0.xml", """<hp:sec xmlns:hp="urn:hp" xmlns:hc="urn:hc"><hp:p><hp:t>균형교양 12학점</hp:t><hc:img binaryItemIDRef="image1"/></hp:p></hp:sec>""")
        zf.writestr("Contents/BinData/image1.png", _png())

    result = extract_hwpx_figures(source.getvalue(), "", tmp_path, _analyzer)

    _assert_contract(result, "hwpx_binary_item_reference")


def test_xlsx_figure_uses_cell_anchor(tmp_path):
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "교양"
    sheet["A1"] = "교육과정"
    sheet["A2"] = "학문기초교양"
    image_stream = io.BytesIO(_png())
    image_stream.name = "filter.png"
    sheet.add_image(XLImage(image_stream), "B2")
    output = io.BytesIO()
    workbook.save(output)

    result = extract_xlsx_figures(output.getvalue(), "[Sheet: 교양]\n[\ud45c \ud5e4\ub354] 교육과정\n[\ud589] 학문기초교양", tmp_path, _analyzer)

    _assert_contract(result, "xlsx_drawing_anchor")
    assert "B2" in result["figures"][0]["context"]


def test_pptx_figure_uses_slide_shape_order(tmp_path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_png()), 0, 0)
    textbox = slide.shapes.add_textbox(0, 100, 500, 100)
    textbox.text = "엑셀 필터 방법"
    output = io.BytesIO()
    presentation.save(output)

    result = extract_pptx_figures(output.getvalue(), "", tmp_path, _analyzer)

    _assert_contract(result, "pptx_slide_shape_order")
    assert "슬라이드 1" in result["figures"][0]["context"]


def test_pdf_figure_uses_page_level_context(tmp_path):
    image = Image.new("RGB", (180, 100), "white")
    source = io.BytesIO()
    image.save(source, "PDF")

    result = extract_pdf_figures(source.getvalue(), "PDF 본문", tmp_path, _analyzer)

    _assert_contract(result, "pdf_page_image_object")
    assert "PDF 1페이지" in result["figures"][0]["context"]
