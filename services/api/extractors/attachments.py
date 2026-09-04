"""공주대 공지 첨부파일/본문 이미지 → 텍스트 변환 어댑터.

각 어댑터는 실패 시 빈 문자열 또는 [실패 사유]를 돌려준다.
호출자는 결과를 본문에 그대로 이어 붙이면 된다.
"""
from model import get_llm, image_to_text
from parsers.pdf_parser import parse_pdf      # ODL 기반 PDF→마크다운 공유 헬퍼
from extractors.hwp_structured import extract_hwp_structured
from extractors.structured_figures import (
    extract_docx_figures,
    extract_hwpx_figures,
    extract_pdf_figures,
    extract_pptx_figures,
    extract_xlsx_figures,
)

# --- 표준 라이브러리 ---
import io           # 바이트 데이터를 "파일처럼" 다루기 위한 BytesIO 용도 (zipfile/이미지 버퍼가 파일객체를 요구함)
import json
import zipfile      # HWPX 파일은 사실상 ZIP 컨테이너라서 직접 열어서 내부 XML을 꺼냄
import os
import re
import struct
import zlib
from pathlib import Path                       # 파일 확장자(.pdf, .hwpx 등) 추출용
from xml.etree import ElementTree as ET        # HWPX 내부 XML 파싱

# --- 외부 라이브러리 ---
import openpyxl                                # XLSX 읽기 (현재 라우터에서는 미사용)
import olefile                                 # HWP 5.x OLE 컨테이너/압축 문단 직접 추출
import xlrd                                    # XLS 읽기
from PIL import Image
from pptx import Presentation


# HWPX 본문(paragraph)의 XML 네임스페이스.
# 이 prefix를 붙여야 ElementTree가 <hp:t> 같은 텍스트 노드를 찾을 수 있다.
_HWPX_PARA_NS = "{http://www.hancom.co.kr/hwpml/2011/paragraph}"
_SUPPORTED_ZIP_EXTS = {".zip", ".pdf", ".docx", ".hwpx", ".hwp", ".xlsx", ".xls", ".ppt", ".pptx", ".jpg", ".jpeg", ".png", ".gif"}
_MAX_ZIP_MEMBERS = 30
_MAX_ZIP_DEPTH = 2
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"}


def hwpx_bytes_to_text(data: bytes) -> str:
    """hwpx 패키지(ZIP)의 Contents/section*.xml에서 <hp:t> 텍스트 노드를 모두 추출.

    표/일정/문의처까지 텍스트로 들어있다. 이미지는 누락되지만 RAG 입력에는 충분.
    """
    parts = []  # 모든 section에서 모은 텍스트 조각을 담을 버퍼

    # bytes를 파일처럼 감싸서 ZipFile에 넘긴다 (디스크에 저장하지 않고 메모리에서 처리)
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        # HWPX 내부 구조: Contents/section0.xml, section1.xml, ... 식으로 본문이 분할 저장돼 있음
        # 페이지 순서를 맞추기 위해 정렬해서 순회한다
        names = sorted(
            n for n in z.namelist()
            if n.startswith("Contents/section") and n.endswith(".xml")
        )
        for name in names:
            # 해당 section XML을 읽어 트리로 파싱
            root = ET.fromstring(z.read(name))
            # <hp:t> = HWPX의 "text run" 노드. 실제 글자가 담긴 leaf 노드만 골라낸다
            for t in root.iter(f"{_HWPX_PARA_NS}t"):
                if t.text:  # 빈 노드는 스킵
                    parts.append(t.text)

    # 줄바꿈으로 이어붙여 단일 문자열로 반환 (앞뒤 공백 정리)
    return "\n".join(parts).strip()


def docx_bytes_to_text(data: bytes) -> str:
    """DOCX OpenXML의 문단·표 텍스트를 문서 순서대로 추출한다."""
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        root = ET.fromstring(archive.read("word/document.xml"))
    paragraphs = []
    for node in root.iter():
        if node.tag.rsplit("}", 1)[-1] != "p":
            continue
        text = "".join(
            str(child.text or "")
            for child in node.iter()
            if child.tag.rsplit("}", 1)[-1] == "t"
        ).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs).strip()


# VLM(Gemini)에게 OCR을 시킬 때 쓰는 고정 프롬프트.
# "설명 문장 붙이지 마라"가 핵심 — 안 그러면 "이 이미지는 ~에 대한 안내입니다" 같은 군더더기가 본문에 섞임.
_VLM_PROMPT = """이 이미지는 대학 공지글의 일부다. 이미지에 적힌 모든 텍스트와 표를 한국어 plain text로 빠짐없이 추출하라.
원칙:
- 설명, 분석, 추론, 생각 과정 출력 절대 금지.
- 행사명, 일정, 신청기한, 신청방법, 문의처 등 정보 항목은 누락 없이 그대로 옮긴다.
- 표는 줄바꿈으로 항목을 구분한다.
- 장식/광고 문구도 모두 포함한다.
- 텍스트만 출력하고 설명 문장은 붙이지 않는다."""


def _image_to_text(image_bytes: bytes, mime: str, prompt: str = _VLM_PROMPT) -> str:
    """이미지 바이트를 Gemini에 던져 텍스트만 받아오는 저수준 헬퍼."""
    return image_to_text(image_bytes, mime, prompt)


def _download(url: str, context) -> bytes:
    """첨부 파일을 HTTP로 직접 받는다(브라우저 context의 쿠키/세션 공유). 실패 시 빈 바이트.

    download.do처럼 Content-Disposition: attachment 응답도 browser navigation 없이
    바이트로 받는다. page.goto는 그런 응답에서 navigation을 abort하며 예외를 던져
    받은 파일까지 버리므로 쓰지 않는다.
    """
    try:
        response = context.request.get(
            url,
            timeout=int(os.getenv("ATTACHMENT_DOWNLOAD_TIMEOUT_SECONDS", "180000")),
            fail_on_status_code=False,
        )
        if response.ok and response.body():
            return response.body()
        print(f"[download failed] {url} -> status={response.status}")
        return b""
    except Exception as e:
        print(f"[download failed] {url} -> {e}")
        return b""


def pdf_to_text(data: bytes) -> str:
    """텍스트 PDF용 1차 추출(ODL 마크다운, 표는 순수 md). 스캔 PDF는 빈 문자열을 반환한다."""
    return parse_pdf(data, markdown_with_html=False)


def _pdf_bytes_full(data: bytes) -> str:
    """ODL 1차(텍스트/표 마크다운) → 비어있으면 pdf2image+VLM fallback."""
    try:
        body = pdf_to_text(data)
    except Exception as e:
        # ODL(JVM) 변환 실패/타임아웃 → 빈값 취급해 아래 VLM 폴백으로 넘긴다.
        print(f"[odl failed] {e}")
        body = ""
    if body:
        # 텍스트 레이어가 있는 정상 PDF: 1차 결과를 그대로 사용
        return body

    # 여기 도달했다는 건 "스캔본 PDF" = 이미지 덩어리. 페이지를 렌더링해서 OCR로 돌린다.
    # pdf2image는 무거운 의존성이라 폴백 경로에서만 lazy import.
    from pdf2image import convert_from_bytes
    images = convert_from_bytes(data, dpi=150)  # 150dpi면 OCR 품질과 속도의 합리적 절충

    chunks = []
    for im in images:
        # PIL 이미지를 PNG 바이트로 직렬화 → VLM에 전달
        buf = io.BytesIO()
        im.save(buf, format="PNG")
        chunks.append(_image_to_text(buf.getvalue(), "image/png"))
    return "\n".join(chunks).strip()


def _hwp_para_payload_to_text(payload: bytes) -> str:
    """HWP 5.x HWPTAG_PARA_TEXT payload에서 제어 문자를 빼고 글자만 복원."""
    units = struct.unpack(f"<{len(payload) // 2}H", payload[: len(payload) // 2 * 2])
    extended_controls = set(range(1, 10)) | set(range(11, 13)) | set(range(14, 24))
    out: list[str] = []
    index = 0
    while index < len(units):
        code = units[index]
        if code in extended_controls:
            if code == 9:
                out.append(" ")
            # 필드/표/그림 등 inline control은 8 UTF-16 code unit 고정 길이다.
            index += 8
            continue
        if code in (10, 13):
            out.append("\n")
        elif code in (24, 30, 31):
            out.append(" ")
        elif 0xD800 <= code <= 0xDBFF:
            if index + 1 < len(units) and 0xDC00 <= units[index + 1] <= 0xDFFF:
                low = units[index + 1]
                out.append(chr(0x10000 + ((code - 0xD800) << 10) + (low - 0xDC00)))
                index += 1
            else:
                out.append("\uFFFD")
        elif 0xDC00 <= code <= 0xDFFF:
            out.append("\uFFFD")
        elif code >= 32:
            out.append(chr(code))
        index += 1

    text = "".join(out).replace("\x00", " ")
    text = re.sub(r"[ \t\u3000]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    return text.strip()


def _hwp_record_stream_text(stream: bytes) -> str:
    """HWP record stream의 HWPTAG_PARA_TEXT(tag 67) 레코드만 순서대로 추출."""
    paragraphs: list[str] = []
    offset = 0
    while offset + 4 <= len(stream):
        header = struct.unpack_from("<I", stream, offset)[0]
        offset += 4
        tag_id = header & 0x3FF
        size = (header >> 20) & 0xFFF
        if size == 0xFFF:
            if offset + 4 > len(stream):
                break
            size = struct.unpack_from("<I", stream, offset)[0]
            offset += 4
        if offset + size > len(stream):
            break
        payload = stream[offset : offset + size]
        offset += size
        if tag_id == 67:
            text = _hwp_para_payload_to_text(payload)
            if text:
                paragraphs.append(text)
    return "\n".join(paragraphs).strip()


def _hwp5_ole_body_text(data: bytes) -> str:
    """HWP 5.x OLE BodyText section을 직접 풀어 LibreOffice 없이 본문 추출."""
    with olefile.OleFileIO(io.BytesIO(data)) as document:
        if not document.exists("FileHeader"):
            return ""
        header = document.openstream("FileHeader").read()
        if len(header) < 40:
            return ""
        compressed = bool(struct.unpack_from("<I", header, 36)[0] & 1)
        section_names = [
            "/".join(parts)
            for parts in document.listdir()
            if len(parts) == 2
            and parts[0] == "BodyText"
            and parts[1].startswith("Section")
        ]
        section_names.sort(key=lambda name: int(name.rsplit("Section", 1)[1]))

        sections: list[str] = []
        for name in section_names:
            raw = document.openstream(name).read()
            try:
                stream = zlib.decompress(raw, -15) if compressed else raw
            except zlib.error:
                continue
            text = _hwp_record_stream_text(stream)
            if text:
                sections.append(text)
        return "\n".join(sections).strip()


def hwp_bytes_to_text(data: bytes, filename: str = "attachment.hwp") -> str:
    """독립적인 HWP 5.x BodyText 레코드 파서로 문단을 추출한다."""
    try:
        text = _hwp5_ole_body_text(data)
    except (OSError, ValueError, struct.error, olefile.OleFileError):
        text = ""
    if text:
        return text

    raise RuntimeError(f"HWP 5.x BodyText 구조 추출 실패: {filename}")


def _zip_member_text(
    member_name: str,
    data: bytes,
    source_url: str,
    context,
    include_xlsx: bool,
    depth: int,
) -> str:
    ext = Path(member_name.lower()).suffix
    label = f"[압축 내부 파일: {member_name}]"

    if ext == ".zip":
        nested = _zip_bytes_to_text(data, source_url, context, include_xlsx, depth + 1)
        return f"{label}\n{nested}" if nested else f"{label}\n(압축 내부 텍스트 없음)"

    if ext == ".pdf":
        return f"{label}\n{_pdf_bytes_full(data)}"

    if ext == ".docx":
        return f"{label}\n{docx_bytes_to_text(data)}"

    if ext == ".pptx":
        return f"{label}\n{pptx_to_text(data)}"

    if ext == ".ppt":
        return f"{label}\n(구형 PPT는 정확한 구조 추출기가 없어 검토 대상으로 보존)"

    if ext == ".hwpx":
        return f"{label}\n{hwpx_bytes_to_text(data)}"

    if ext == ".hwp":
        return f"{label}\n{hwp_bytes_to_text(data, member_name)}"

    if ext == ".xlsx":
        if include_xlsx:
            return f"{label}\n{xlsx_to_text(data)}"
        return f"{label}\n(엑셀 첨부 — 임베딩 제외. 원본 ZIP: {source_url})"

    if ext == ".xls":
        if include_xlsx:
            return f"{label}\n{xls_to_text(data)}"
        return f"{label}\n(엑셀 첨부 — 임베딩 제외. 원본 ZIP: {source_url})"

    if ext in _IMAGE_EXTS:
        mime = "image/png" if ext == ".png" else "image/jpeg"
        return f"{label}\n{_image_to_text(data, mime).strip()}"

    return f"{label}\n(지원하지 않는 압축 내부 확장자, 건너뜀)"


def _zip_bytes_to_text(data: bytes, source_url: str, context, include_xlsx: bool, depth: int = 0) -> str:
    if depth > _MAX_ZIP_DEPTH:
        return "(ZIP 중첩 깊이 제한으로 내부 압축 해제 중단)"

    parts: list[str] = []
    if not data:
        return "(빈 ZIP 데이터)"
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        members = [
            info for info in z.infolist()
            if not info.is_dir()
            and "__MACOSX/" not in info.filename
            and not Path(info.filename).name.startswith(".")
        ]
        handled = 0
        for info in members:
            if handled >= _MAX_ZIP_MEMBERS:
                parts.append(f"(ZIP 내부 파일 {len(members)}개 중 {_MAX_ZIP_MEMBERS}개만 처리)")
                break

            ext = Path(info.filename.lower()).suffix
            if ext not in _SUPPORTED_ZIP_EXTS:
                continue

            handled += 1
            try:
                parts.append(
                    _zip_member_text(
                        info.filename,
                        z.read(info),
                        source_url,
                        context,
                        include_xlsx,
                        depth,
                    )
                )
            except Exception as e:
                parts.append(f"[압축 내부 파일: {info.filename}]\n(처리 실패: {e})")

    return "\n\n".join(part for part in parts if part.strip()).strip()


# XLSX 첨부는 노이즈가 많아 기본은 임베딩 제외. 단 아래 키워드가 제목/본문에 보이면 활성화.
# 수강신청·교양·교과목 편성 같은 표 데이터가 핵심인 공지에서만 켜진다.
XLSX_KEYWORDS = ("교양", "수강신청", "교과목", "편성", "시간표", "강의계획", "개설")


def _trim_empty_tail(cells: list[str]) -> list[str]:
    """오른쪽 끝 빈 셀 제거. 엑셀은 사용 범위가 넓게 잡히는 일이 많다."""
    end = len(cells)
    while end > 0 and not cells[end - 1].strip():
        end -= 1
    return cells[:end]


def _looks_numeric_or_code(value: str) -> bool:
    s = value.strip().replace(",", "")
    if not s:
        return False
    if s.replace(".", "", 1).isdigit():
        return True
    # 2026-1, 3-3-0, 2009094 같은 학기/학점/코드형 값.
    compact = s.replace("-", "").replace(".", "").replace("/", "")
    return compact.isdigit()


def _looks_like_header(row: list[str], next_rows: list[list[str]]) -> bool:
    non_empty = [c.strip() for c in row if c.strip()]
    if len(non_empty) < 2:
        return False

    following_width = max(
        (sum(1 for c in r if c.strip()) for r in next_rows),
        default=0,
    )
    if following_width < 2:
        return False

    unique_ratio = len(set(non_empty)) / len(non_empty)
    textish_ratio = sum(not _looks_numeric_or_code(c) for c in non_empty) / len(non_empty)
    avg_len = sum(len(c) for c in non_empty) / len(non_empty)

    return unique_ratio >= 0.7 and textish_ratio >= 0.55 and avg_len <= 40


def _dedupe_headers(row: list[str]) -> list[str]:
    headers = []
    seen: dict[str, int] = {}
    for idx, cell in enumerate(row):
        base = cell.strip() or f"열{idx + 1}"
        seen[base] = seen.get(base, 0) + 1
        headers.append(base if seen[base] == 1 else f"{base}_{seen[base]}")
    return headers


def pptx_to_text(data: bytes) -> str:
    """PPTX의 텍스트 프레임과 표를 슬라이드 순서대로 추출한다."""
    presentation = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for slide_no, slide in enumerate(presentation.slides, start=1):
        parts = [f"[슬라이드 {slide_no}]"]
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        chunks.append("\n".join(parts))
    return "\n\n".join(chunks).strip()


_HWP_IMAGE_PROMPT = """이 이미지는 대학 공지 HTML 본문 또는 문서(HWP·HWPX·PDF·Word·Excel·PowerPoint)에서 추출한 원본 이미지다.
반드시 아래 JSON 객체 하나만 출력하라.
{"kind":"music_score|scanned_text|table_image|map|chart|photo|logo|other","ocrText":"보이는 글자를 원문 그대로","description":"문서 문맥을 반영한 짧은 설명","contextMatch":"supports|unrelated|uncertain","confidence":0.0,"music":{"title":"","lyricist":"","composer":"","lyrics":"","notationStatus":"preserved_as_image"}}
악보이면 제목·작사·작곡·가사처럼 눈으로 확실히 읽히는 정보만 기록한다. 음표나 음정을 추측하거나 악보를 텍스트 음계로 만들어내지 말고 notationStatus는 preserved_as_image로 둔다.
일반 이미지도 보이지 않는 내용을 추측하지 않는다. 주어진 앞뒤 문맥과 이미지가 실제로 서로를 보완하는지 contextMatch로 판정하고, 문맥과 충돌하면 unrelated 또는 uncertain으로 둔다."""


def _decode_image_analysis_json(raw: str) -> dict | None:
    """작은 VLM이 붙이는 fence·잘못된 역슬래시·trailing comma를 제한적으로 복구."""
    candidates = [raw.strip()]
    fenced = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw, flags=re.I | re.S)
    if fenced:
        candidates.insert(0, fenced.group(1).strip())
    first = raw.find("{")
    last = raw.rfind("}")
    if first >= 0 and last > first:
        candidates.append(raw[first:last + 1].strip())

    for candidate in candidates:
        variants = [candidate]
        repaired = re.sub(r'\\(?!["\\/bfnrtu])', "", candidate)
        repaired = re.sub(r",\s*([}\]])", r"\1", repaired)
        if repaired != candidate:
            variants.append(repaired)
        for variant in variants:
            try:
                value = json.loads(variant, strict=False)
            except json.JSONDecodeError:
                continue
            if isinstance(value, dict):
                return value
    return None


def _hwp_image_analysis(image_bytes: bytes, mime: str, filename: str, document_context: str = "") -> dict:
    """문서 내부 이미지를 OCR/VLM으로 분류하되 악보 음표는 추측하지 않는다."""
    raw = _image_to_text(
        image_bytes,
        mime,
        f"파일명: {filename}\n문서 안의 그림 위치와 앞뒤 문맥:\n{document_context or '(문맥 없음)'}\n\n{_HWP_IMAGE_PROMPT}",
    ).strip()
    value = _decode_image_analysis_json(raw)
    if value is None:
        return {
            "kind": "other",
            "ocrText": "",
            "description": re.sub(r"```(?:json)?|```", "", raw, flags=re.I).strip()[:1200],
            "confidence": 0.0,
            "status": "unstructured_response",
        }
    allowed = {"music_score", "scanned_text", "table_image", "map", "chart", "photo", "logo", "other"}
    if value.get("kind") not in allowed:
        value["kind"] = "other"
    if value.get("contextMatch") not in {"supports", "unrelated", "uncertain"}:
        value["contextMatch"] = "uncertain"
    if value["kind"] == "music_score":
        music = value.get("music") if isinstance(value.get("music"), dict) else {}
        music["notationStatus"] = "preserved_as_image"
        value["music"] = music
    return value


def _figure_analyzer():
    enabled = os.getenv("DOCUMENT_IMAGE_ANALYSIS_ENABLED")
    if enabled is None:
        enabled = os.getenv("HWP_IMAGE_ANALYSIS_ENABLED", "false")
    return (
        _hwp_image_analysis
        if enabled.lower() in {"1", "true", "yes", "on"}
        else None
    )


def _document_assets_root() -> Path:
    return Path(
        os.getenv("DOCUMENT_ASSETS_ROOT")
        or os.getenv("HWP_ASSETS_ROOT", "data/assets")
    )


def _attach_structured_figures(meta: dict, structured: dict) -> str:
    meta["derived_assets"] = list(meta.get("derived_assets") or []) + list(
        structured.get("derived_assets") or []
    )
    meta["figure_contents"] = list(meta.get("figure_contents") or []) + list(
        structured.get("figure_contents") or []
    )
    meta["figure_bundle_dir"] = structured.get("bundle_dir")
    return str(structured.get("text") or "")


def xlsx_relevant(*texts: str) -> bool:
    """제목·본문 등을 합쳐 XLSX_KEYWORDS 중 하나라도 포함하면 True."""
    blob = "\n".join(t for t in texts if t)
    return any(kw in blob for kw in XLSX_KEYWORDS)


def xlsx_to_text(data: bytes) -> str:
    """XLSX 시트를 검색 친화적인 텍스트로 직렬화.

    표 헤더는 [표 헤더]로 표시하고, 데이터 행은 [행] 값 나열로 둔다.
    청킹 단계(embed.py)가 표 안에서 잘린 chunk 앞에 현재 헤더를 다시
    붙여주므로, 여기서는 행마다 컬럼명을 반복하지 않는다.
    """
    # data_only=True: 수식(=A1+B1) 대신 마지막으로 계산된 값을 가져옴
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    out = []
    for sheet in wb.worksheets:
        out.append(f"[Sheet: {sheet.title}]")  # 시트 경계 표시
        headers: list[str] | None = None
        rows = []
        for row in sheet.iter_rows(values_only=True):
            # None은 빈 문자열로, 나머지는 str로 일괄 변환
            cells = _trim_empty_tail(["" if v is None else str(v).strip() for v in row])
            # 완전히 빈 행은 스킵 (엑셀에 흔한 공백 행 제거)
            if any(c.strip() for c in cells):
                rows.append(cells)

        for idx, row in enumerate(rows):
            if headers is None and _looks_like_header(row, rows[idx + 1:idx + 4]):
                headers = _dedupe_headers(row)
                out.append("[표 헤더] " + " | ".join(headers))
                continue

            if headers:
                row_text = " | ".join(c for c in row if c)
                if not row_text:
                    continue
                out.append("[행] " + row_text)
            else:
                out.append(" | ".join(c for c in row if c))
        out.append(f"[End Sheet: {sheet.title}]")
    return "\n".join(out).strip()


def xls_to_text(data: bytes) -> str:
    """구형 XLS 시트를 검색 친화적인 텍스트로 직렬화."""
    book = xlrd.open_workbook(file_contents=data)
    out = []
    for sheet in book.sheets():
        out.append(f"[Sheet: {sheet.name}]")
        headers: list[str] | None = None
        rows = []
        for r in range(sheet.nrows):
            cells = _trim_empty_tail([
                "" if sheet.cell_value(r, c) is None else str(sheet.cell_value(r, c)).strip()
                for c in range(sheet.ncols)
            ])
            if any(c.strip() for c in cells):
                rows.append(cells)

        for idx, row in enumerate(rows):
            if headers is None and _looks_like_header(row, rows[idx + 1:idx + 4]):
                headers = _dedupe_headers(row)
                out.append("[표 헤더] " + " | ".join(headers))
                continue

            row_text = " | ".join(c for c in row if c)
            if row_text:
                out.append(("[행] " if headers else "") + row_text)
        out.append(f"[End Sheet: {sheet.name}]")
    return "\n".join(out).strip()


def _office_preview_fallback(att: dict, context) -> str:
    preview_url = att.get("preview_url")

    if not preview_url:
        return ""

    try:
        text = _preview_via_browser(preview_url, context)

        # viewer에서 일부 텍스트라도 확보되면 그대로 사용한다.
        # 공주대처럼 download.do는 막혀 있지만 synapView 렌더는 허용하는
        # 사이트가 있어 preview 텍스트를 우선 신뢰한다.
        return text or ""

    except Exception as e:
        print(f"[preview fallback failed] {preview_url} -> {e}")
        return ""


def _preview_via_browser(preview_url: str, context) -> str:
    """구조 파서가 없는 레거시 첨부에만 쓰는 synapView 최후 폴백.

    synapView는 대형 HWP/HWPX를 lazy rendering 하는 경우가 많아
    단순 body 추출만으로는 앞부분 일부만 수집된다.
    따라서 viewer를 끝까지 스크롤하며 iframe 텍스트 변화를 반복 수집한다.
    """
    page = context.new_page()

    try:
        try:
            page.goto(
                preview_url,
                wait_until="domcontentloaded",
                timeout=int(os.getenv("HWP_PREVIEW_TIMEOUT_MS", "45000")),
            )
        except Exception as e:
            # synapView는 timeout 이후에도 iframe 렌더가 계속 진행되는 경우가 많다.
            # 따라서 timeout을 치명적 실패로 보지 않고 계속 스크롤 수집을 시도한다.
            print(f"[synap goto timeout ignored] {e}")

        page.wait_for_timeout(
            int(os.getenv("HWP_PREVIEW_EXTRA_WAIT_MS", "10000"))
        )

        max_scrolls = int(os.getenv("HWP_PREVIEW_MAX_SCROLLS", "1200"))
        scroll_px = int(os.getenv("HWP_PREVIEW_SCROLL_PX", "4000"))
        settle_ms = int(os.getenv("HWP_PREVIEW_SCROLL_WAIT_MS", "3500"))
        stable_limit = int(os.getenv("HWP_PREVIEW_STABLE_LIMIT", "30"))

        collected_chunks: list[str] = []
        seen_texts: set[str] = set()

        stable_count = 0
        viewer_selectors = [
            ".viewer",
            "#viewer",
            ".doc-view",
            ".synap-viewer",
            ".document-view",
            ".viewer-container",
        ]
        last_total_len = 0

        def collect_frame_texts() -> int:
            total = 0

            for frame in page.frames:
                try:
                    text = frame.inner_text("body").strip()
                except Exception:
                    continue

                if not text:
                    continue

                normalized = "\n".join(
                    line.strip()
                    for line in text.splitlines()
                    if line.strip()
                )

                if not normalized:
                    continue

                total += len(normalized)

                if normalized not in seen_texts:
                    seen_texts.add(normalized)
                    collected_chunks.append(normalized)

            return total

        # 초기 수집
        last_total_len = collect_frame_texts()

        for _ in range(max_scrolls):
            scrolled = False

            for selector in viewer_selectors:
                try:
                    locator = page.locator(selector).first

                    if locator.count() == 0:
                        continue

                    locator.evaluate(
                        "(el, y) => el.scrollBy(0, y)",
                        scroll_px,
                    )

                    scrolled = True
                    break

                except Exception:
                    continue

            # viewer container를 못 찾으면 body wheel fallback
            if not scrolled:
                try:
                    page.mouse.wheel(0, scroll_px)
                except Exception:
                    pass

            page.wait_for_timeout(settle_ms)

            current_total_len = collect_frame_texts()

            # 더 이상 텍스트 증가가 없으면 안정화 카운트 증가
            if current_total_len <= last_total_len:
                stable_count += 1
            else:
                stable_count = 0
                last_total_len = current_total_len

            # 여러 번 연속 변화 없으면 종료
            if stable_count >= stable_limit:
                break

        print(
            f"[synap collected] chunks={len(collected_chunks)} chars={sum(len(c) for c in collected_chunks)}"
        )
        merged = "\n\n".join(collected_chunks).strip()

        # 너무 긴 중복 제거
        lines = []
        seen_lines = set()

        for line in merged.splitlines():
            normalized = line.strip()

            if not normalized:
                continue

            if normalized in seen_lines:
                continue

            seen_lines.add(normalized)
            lines.append(normalized)

        return "\n".join(lines).strip()

    finally:
        try:
            page.close()
        except Exception:
            pass


def inline_image_to_text(
    image_url: str,
    context,
    document_context: str = "",
    filename: str = "notice-body-image",
) -> tuple[str, bytes | None, str | None, dict]:
    """본문 inline 이미지를 문맥과 함께 VLM으로 구조화한다.

    반환: (search_text, raw_bytes, mime, analysis). 작은 UI 이미지와 로고는
    검색·렌더링 자산으로 채택하지 않는다.
    """
    try:
        data = _download(image_url, context)
    except Exception:
        return "", None, None, {"status": "download_failed"}
    if not data:
        return "", None, None, {"status": "download_failed"}

    mime_by_format = {
        "PNG": "image/png", "JPEG": "image/jpeg", "GIF": "image/gif",
        "WEBP": "image/webp", "BMP": "image/bmp",
    }
    width = height = None
    try:
        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
            mime = mime_by_format.get(str(image.format or "").upper(), "image/jpeg")
    except Exception:
        mime = "image/png" if image_url.lower().split("?", 1)[0].endswith(".png") else "image/jpeg"

    if width is not None and height is not None and (
        width < 48 or height < 32 or width * height < 4096 or max(width, height) > 20 * max(1, min(width, height))
    ):
        return "", None, None, {
            "status": "ignored_small_ui_image", "width": width, "height": height,
        }

    try:
        analysis = _hwp_image_analysis(data, mime, filename, document_context)
    except Exception as error:
        analysis = {"status": "failed", "error": f"{type(error).__name__}: {error}"}
    analysis["width"] = width
    analysis["height"] = height
    search_text = "\n".join(value for value in (
        str(analysis.get("ocrText") or "").strip(),
        str(analysis.get("description") or "").strip(),
    ) if value)
    analysis["searchText"] = search_text
    analysis["documentContext"] = document_context
    if analysis.get("kind") == "logo":
        return "", None, None, {**analysis, "status": "ignored_logo"}
    return search_text, data, mime, analysis


def attachment_to_text(att: dict, context, include_xlsx: bool = False):
    """att = {'filename', 'download_url', 'preview_url' | None}.

    include_xlsx: 기본 False. True면 XLSX 본문도 추출해 임베딩 대상에 포함.
                  호출자(crawler)가 xlsx_relevant(title, body)로 판정해 전달.

    반환: (text, asset_meta)
      text: '[첨부: <파일명>]\\n<본문>' (기존 호환, 실패 시 본문 자리에 사유)
      asset_meta: {
        kind: inline_image | attachment_image | attachment_pdf
              | attachment_hwpx | attachment_xlsx | attachment_other,
        filename, source_url, mime_type,
        raw_bytes: bytes | None,  — attachment_image일 때만 채워짐
        extracted_text: str
      }
    """
    # 입력 unpack
    name = att["filename"]
    ext = Path(name.lower()).suffix       # .pdf / .hwpx / .jpg ... — 소문자 통일 후 확장자 추출
    label = f"[첨부: {name}]"             # 본문 앞에 붙일 라벨 (RAG 컨텍스트에서 출처 식별용)
    source_url = att["download_url"]

    if not source_url:
        body = "(첨부 다운로드 URL 없음)"
        meta = {
            "kind": "attachment_other",
            "filename": name,
            "source_url": source_url,
            "mime_type": None,
            "raw_bytes": None,
            "extracted_text": body,
        }
        return f"{label}\n{body}", meta

    # 메타 기본값을 먼저 깔아두고, 아래 분기에서 필드를 덮어쓰는 패턴
    meta = {
        "kind": "attachment_other",
        "filename": name,
        "source_url": source_url,
        "mime_type": None,
        "raw_bytes": None,
        "extracted_text": "",
    }

    # try 전체로 감싸서 어떤 분기에서 예외가 나더라도 "(처리 실패: ...)" 문자열로 환원한다
    # → 호출자(crawler)는 예외 처리 없이 결과를 본문에 그대로 이어붙일 수 있다
    try:
        # ───────── 분기 1: PDF ─────────
        if ext == ".pdf":
            meta["kind"] = "attachment_pdf"
            meta["mime_type"] = "application/pdf"
            data = _download(source_url, context)
            meta["raw_bytes"] = data
            body = _pdf_bytes_full(data)   # 텍스트 1차 → 실패 시 이미지 OCR 폴백 (위 함수 참고)
            try:
                body = _attach_structured_figures(
                    meta,
                    extract_pdf_figures(
                        data, body, _document_assets_root(), _figure_analyzer()
                    ),
                )
            except Exception as error:
                meta["figure_extraction_error"] = f"{type(error).__name__}: {error}"

        # ───────── 분기 1.25: Word OpenXML ─────────
        elif ext == ".docx":
            meta["kind"] = "attachment_docx"
            meta["mime_type"] = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            data = _download(source_url, context)
            meta["raw_bytes"] = data
            structured = extract_docx_figures(
                data, "", _document_assets_root(), _figure_analyzer()
            )
            body = _attach_structured_figures(meta, structured)

        # ───────── 분기 1.5: PPT / PPTX ─────────
        elif ext in (".ppt", ".pptx"):
            meta["kind"] = "attachment_ppt"
            meta["mime_type"] = (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                if ext == ".pptx" else "application/vnd.ms-powerpoint"
            )
            data = _download(source_url, context)
            meta["raw_bytes"] = data
            if ext == ".pptx":
                body = pptx_to_text(data)
                try:
                    body = _attach_structured_figures(
                        meta,
                        extract_pptx_figures(
                            data, body, _document_assets_root(), _figure_analyzer()
                        ),
                    )
                except Exception as error:
                    meta["figure_extraction_error"] = f"{type(error).__name__}: {error}"
            else:
                body = "(구형 PPT는 정확한 구조 추출기가 없어 검토 대상으로 보존)"
                meta["review_required"] = True
                meta["review_reason"] = "unsupported_legacy_ppt"

        # ───────── 분기 2: 엑셀 ─────────
        elif ext in (".xlsx", ".xls"):
            meta["kind"] = "attachment_xlsx"
            meta["mime_type"] = (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                if ext == ".xlsx" else "application/vnd.ms-excel"
            )
            if include_xlsx and ext == ".xlsx":
                # 키워드 매칭(수강신청·교양·편성 등)이 걸린 공지 → 표 전체를 텍스트화해서 임베딩 대상에 포함.
                try:
                    data = _download(source_url, context)
                    meta["raw_bytes"] = data
                    body = xlsx_to_text(data)
                    body = _attach_structured_figures(
                        meta,
                        extract_xlsx_figures(
                            data, body, _document_assets_root(), _figure_analyzer()
                        ),
                    )
                except Exception:
                    body = _office_preview_fallback(att, context)
                    if not body:
                        raise
            elif include_xlsx and ext == ".xls":
                try:
                    data = _download(source_url, context)
                    meta["raw_bytes"] = data
                    body = xls_to_text(data)
                except Exception:
                    body = _office_preview_fallback(att, context)
                    if not body:
                        raise
            else:
                # 기본: 노이즈 많은 엑셀은 임베딩 제외하고 안내문만 남김.
                body = f"(엑셀 첨부 — 임베딩 제외. 원본 다운로드: {source_url})"

        # ───────── 분기 3: HWPX / HWP ─────────
        elif ext in (".hwpx", ".hwp"):
            meta["kind"] = "attachment_hwpx" if ext == ".hwpx" else "attachment_hwp"
            meta["mime_type"] = (
                "application/vnd.hancom.hwpx" if ext == ".hwpx" else "application/x-hwp"
            )
            if ext == ".hwp":
                file_data = _download(source_url, context)
                meta["raw_bytes"] = file_data
                if not file_data:
                    body = "(원본 HWP 다운로드 실패)"
                    meta["review_required"] = True
                    meta["review_reason"] = "hwp_download_failed"
                else:
                    secondary_text = hwp_bytes_to_text(file_data, name)
                    structured = extract_hwp_structured(
                        file_data,
                        name,
                        secondary_text=secondary_text,
                        assets_root=_document_assets_root(),
                        image_analyzer=_figure_analyzer(),
                    )
                    derived_assets = [
                        {
                            "kind": "attachment_hwp_structure",
                            "filename": "document.json",
                            "storage_path": structured["structure_path"],
                            "mime_type": "application/json",
                            "extracted_text": json.dumps(
                                structured["quality"], ensure_ascii=False
                            ),
                        },
                        {
                            "kind": "attachment_hwp_markdown",
                            "filename": "document.md",
                            "storage_path": structured["markdown_path"],
                            "mime_type": "text/markdown",
                            "extracted_text": "",
                        },
                    ]
                    if structured.get("converted_hwpx_path"):
                        derived_assets.append({
                            "kind": "attachment_hwp_validation",
                            "filename": "validation.hwpx",
                            "storage_path": structured["converted_hwpx_path"],
                            "mime_type": "application/vnd.hancom.hwpx",
                            "extracted_text": "",
                        })
                    derived_assets.extend(structured["binary_assets"])
                    body = structured["markdown"]
                    meta.update({
                        "storage_path": structured["original_path"],
                        "bundle_dir": structured["bundle_dir"],
                        "markdown_path": structured["markdown_path"],
                        "structure_path": structured["structure_path"],
                        "quality": structured["quality"],
                        "review_required": structured["review_required"],
                        "review_reason": (
                            "hwp_cross_validation_below_threshold"
                            if structured["review_required"] else None
                        ),
                        "derived_assets": derived_assets,
                        "figure_contents": structured.get("figure_contents") or [],
                    })
            else:
                file_data = _download(source_url, context)
                meta["raw_bytes"] = file_data
                if not file_data:
                    body = "(원본 HWPX 다운로드 실패)"
                else:
                    body = hwpx_bytes_to_text(file_data)
                    if not body or not body.strip():
                        raise RuntimeError("HWPX 구조 XML에서 텍스트를 추출하지 못함")
                    body = _attach_structured_figures(
                        meta,
                        extract_hwpx_figures(
                            file_data, body, _document_assets_root(), _figure_analyzer()
                        ),
                    )

        # ───────── 분기 4: 이미지 첨부 ─────────
        elif ext in _IMAGE_EXTS:
            meta["kind"] = "attachment_image"
            meta["mime_type"] = "image/png" if ext == ".png" else "image/jpeg"
            data = _download(source_url, context)
            meta["raw_bytes"] = data       # 멀티모달 임베딩/재처리를 위해 원본 바이트도 보존
            body = _image_to_text(data, meta["mime_type"]).strip()

        # ───────── 분기 5: ZIP ─────────
        elif ext == ".zip":
            meta["kind"] = "attachment_zip"
            meta["mime_type"] = "application/zip"
            data = _download(source_url, context)
            meta["raw_bytes"] = data
            body = _zip_bytes_to_text(data, source_url, context, include_xlsx)
            if not body:
                body = "(ZIP 내부에서 처리 가능한 파일을 찾지 못했습니다.)"

        # ───────── 분기 6: 그 외 확장자 ─────────
        else:
            # 모르는 포맷은 처리 시도조차 하지 않고 안내문만 남기고 early return
            body = "(지원하지 않는 확장자, 건너뜀)"
            meta["extracted_text"] = body
            return f"{label}\n{body}", meta

    except Exception as e:
        # 모듈 docstring의 약속: 예외를 던지지 않고 "(처리 실패: ...)" 문자열로 회수
        body = f"(처리 실패: {e})"
        if ext in {".hwp", ".hwpx", ".ppt"}:
            meta["review_required"] = True
            meta["review_reason"] = f"structured_extraction_failed:{type(e).__name__}"
        meta["extracted_text"] = body
        return f"{label}\n{body}", meta

    # 정상 분기(PDF / 엑셀 / HWPX 성공 / 이미지)의 공통 마무리:
    #   본문이 비어있어도 라벨 + 안내문은 보장해서 호출자가 항상 "라벨\n본문" 형태를 받게 한다
    text = f"{label}\n{body}" if body else f"{label}\n(추출 텍스트 없음)"
    meta["extracted_text"] = body
    return text, meta


def extract_attachment_text(path: str | Path) -> str:
    """파일 경로 기반 generic attachment text extractor.

    curriculum/document crawler 등에서 공통 사용.
    """
    file_path = Path(path)
    ext = file_path.suffix.lower()
    data = file_path.read_bytes()

    if ext == ".pdf":
        return _pdf_bytes_full(data)

    if ext == ".docx":
        return docx_bytes_to_text(data)

    if ext == ".pptx":
        return pptx_to_text(data)

    if ext == ".ppt":
        raise ValueError("구형 PPT는 정확한 구조 추출기가 없어 지원하지 않습니다")

    if ext == ".hwpx":
        return hwpx_bytes_to_text(data)

    if ext == ".hwp":
        return hwp_bytes_to_text(data, file_path.name)

    if ext == ".xlsx":
        return xlsx_to_text(data)

    if ext == ".xls":
        return xls_to_text(data)

    if ext in _IMAGE_EXTS:
        mime = "image/png" if ext == ".png" else "image/jpeg"
        return _image_to_text(data, mime).strip()

    raise ValueError(f"지원하지 않는 attachment 확장자: {ext}")
