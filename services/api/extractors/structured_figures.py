"""Embedded figure extraction for non-HWP document formats.

Every adapter returns the same figure contract used by HWP: numbered placement,
nearby text, a persisted raster asset, optional VLM analysis, inline description,
and an independent retrieval unit.
"""
from __future__ import annotations

import hashlib
import io
import re
import zipfile
from pathlib import Path
from typing import Callable
from xml.etree import ElementTree as ET

import openpyxl
from PIL import Image, ImageOps
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ImageAnalyzer = Callable[[bytes, str, str, str], dict]


def _local(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _image_type(data: bytes, fallback: str = "") -> tuple[str, str] | None:
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return ".png", "image/png"
    if data.startswith(b"\xff\xd8\xff"):
        return ".jpg", "image/jpeg"
    if data.startswith(b"BM"):
        return ".bmp", "image/bmp"
    if data.startswith((b"GIF87a", b"GIF89a")):
        return ".gif", "image/gif"
    suffix = Path(fallback).suffix.lower()
    return ({
        ".png": (".png", "image/png"),
        ".jpg": (".jpg", "image/jpeg"),
        ".jpeg": (".jpg", "image/jpeg"),
        ".bmp": (".bmp", "image/bmp"),
        ".gif": (".gif", "image/gif"),
    }).get(suffix)


def _clip(value: str, limit: int = 900) -> str:
    value = re.sub(r"\s+", " ", value or "").strip()
    return value if len(value) <= limit else value[:limit].rstrip() + "…"


def _context(blocks: list[dict], index: int, location: str) -> str:
    before = next((str(blocks[i].get("text") or "").strip() for i in range(index - 1, -1, -1) if str(blocks[i].get("text") or "").strip()), "")
    after = next((str(blocks[i].get("text") or "").strip() for i in range(index + 1, len(blocks)) if str(blocks[i].get("text") or "").strip()), "")
    return "\n".join(part for part in (
        f"[앞 문맥] {_clip(before, 400)}" if before else "",
        f"[그림 위치] {_clip(location, 500)}" if location else "",
        f"[뒤 문맥] {_clip(after, 400)}" if after else "",
    ) if part)


def _analysis_preview(data: bytes) -> tuple[bytes, str, int, int]:
    with Image.open(io.BytesIO(data)) as opened:
        image = ImageOps.exif_transpose(opened).convert("RGB")
        width, height = image.size
        image.thumbnail((2048, 2048), Image.Resampling.LANCZOS)
        output = io.BytesIO()
        image.save(output, "JPEG", quality=92, optimize=True)
        return output.getvalue(), "image/jpeg", width, height


def _inline_description(analysis: dict) -> str:
    if analysis.get("requiresReview"):
        return ""
    description = str(analysis.get("description") or "").strip()
    ocr = str(analysis.get("ocrText") or "").strip()
    lines = []
    if description:
        lines.append(f"[그림 설명] {description}")
    if ocr and ocr != description:
        lines.append(f"[그림 내 텍스트] {ocr}")
    return "\n".join(lines)


def _finalize(
    *, data: bytes, filename: str, format_name: str, base_text: str,
    blocks: list[dict], assets_root: Path, image_analyzer: ImageAnalyzer | None,
) -> dict:
    digest = hashlib.sha256(data).hexdigest()
    bundle = assets_root / "documents" / digest
    image_dir = bundle / "images"
    image_dir.mkdir(parents=True, exist_ok=True)
    figures: list[dict] = []
    derived_assets: list[dict] = []
    text = base_text
    offset = 0

    for block_index, block in enumerate(blocks):
        for raw in block.get("images") or []:
            image_data = raw.get("data") or b""
            kind = _image_type(image_data, str(raw.get("name") or ""))
            if not image_data or not kind:
                continue
            suffix, mime = kind
            number = len(figures) + 1
            location = str(raw.get("location") or block.get("text") or f"{format_name} block {block_index + 1}")
            context = _context(blocks, block_index, location)
            sha = hashlib.sha256(image_data).hexdigest()
            output_name = f"figure-{number:03d}-{sha[:12]}{suffix}"
            output_path = image_dir / output_name
            output_path.write_bytes(image_data)
            analysis: dict = {}
            width = height = None
            try:
                preview, preview_mime, width, height = _analysis_preview(image_data)
                if image_analyzer and width >= 96 and height >= 40 and width * height >= 12_000:
                    analysis = image_analyzer(preview, preview_mime, output_name, context)
                    analysis["searchText"] = "\n".join(value for value in (
                        str(analysis.get("ocrText") or "").strip(),
                        str(analysis.get("description") or "").strip(),
                    ) if value)
                    analysis["figureNumber"] = number
                    analysis["documentContext"] = context
            except Exception as error:
                analysis = {"status": "failed", "error": f"{type(error).__name__}: {error}"}
            figure = {
                "number": number,
                "label": f"그림 {number}",
                "blockIndex": block_index,
                "context": context,
                "matchMethod": raw.get("matchMethod") or f"{format_name}_document_order",
                "matchConfidence": float(raw.get("matchConfidence", 0.9)),
                "filename": output_name,
                "storagePath": str(output_path),
                "mimeType": mime,
                "sha256": sha,
                "width": width,
                "height": height,
                "analysis": analysis,
            }
            figures.append(figure)
            derived_assets.append({
                "kind": "attachment_document_image",
                "filename": output_name,
                "storage_path": str(output_path),
                "mime_type": mime,
                "sha256": sha,
                "width": width,
                "height": height,
                "analysis": analysis,
                "figure": figure,
            })
            marker = str(raw.get("marker") or f"[그림 {number}]")
            description = _inline_description(analysis)
            replacement = marker + (f"\n{description}" if description else "")
            placeholder = str(raw.get("placeholder") or "")
            if placeholder and placeholder in text:
                pos = text.find(placeholder, offset)
                if pos >= 0:
                    text = text[:pos] + replacement + text[pos + len(placeholder):]
                    offset = pos + len(replacement)
                    continue
            text += f"\n\n{replacement}\n{context}"

    figure_contents = []
    for figure in figures:
        search_text = str(figure["analysis"].get("searchText") or "").strip()
        if not search_text or figure["analysis"].get("requiresReview"):
            continue
        figure_contents.append({
            "number": figure["number"],
            "filename": figure["filename"],
            "text": f"[그림 {figure['number']}]\n{figure['context']}\n[그림 설명] {search_text}",
        })
    return {"text": text.strip(), "figures": figures, "derived_assets": derived_assets, "figure_contents": figure_contents, "bundle_dir": str(bundle)}


def _zip_relations(zf: zipfile.ZipFile, rel_path: str, base_dir: str) -> dict[str, str]:
    if rel_path not in zf.namelist():
        return {}
    root = ET.fromstring(zf.read(rel_path))
    result = {}
    for rel in root:
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target", "")
        if rel_id and target and not target.startswith(("http://", "https://")):
            parts = []
            for part in (Path(base_dir) / target).parts:
                if part == "..":
                    if parts: parts.pop()
                elif part != ".": parts.append(part)
            result[rel_id] = "/".join(parts)
    return result


def extract_docx_figures(data: bytes, base_text: str, assets_root: Path, image_analyzer=None) -> dict:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        root = ET.fromstring(zf.read("word/document.xml"))
        rels = _zip_relations(zf, "word/_rels/document.xml.rels", "word")
        blocks = []
        counter = 0
        body = next((node for node in root.iter() if _local(node.tag) == "body"), root)
        for child in list(body):
            text = " ".join((node.text or "").strip() for node in child.iter() if _local(node.tag) == "t" and (node.text or "").strip())
            images = []
            for blip in (node for node in child.iter() if _local(node.tag) == "blip"):
                rel_id = next((value for key, value in blip.attrib.items() if _local(key) == "embed"), None)
                target = rels.get(rel_id or "")
                if target and target in zf.namelist():
                    counter += 1
                    placeholder = f"[[DOCX_FIGURE_{counter}]]"
                    images.append({"data": zf.read(target), "name": target, "placeholder": placeholder, "location": text or f"문단 {len(blocks)+1}", "matchMethod": "docx_paragraph_relationship", "matchConfidence": 1.0})
                    text = f"{text}\n{placeholder}".strip()
            if text or images: blocks.append({"text": text, "images": images})
    generated = "\n\n".join(block["text"] for block in blocks)
    return _finalize(data=data, filename="document.docx", format_name="docx", base_text=generated or base_text, blocks=blocks, assets_root=assets_root, image_analyzer=image_analyzer)


def extract_hwpx_figures(data: bytes, base_text: str, assets_root: Path, image_analyzer=None) -> dict:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        manifest: dict[str, str] = {}
        if "Contents/content.hpf" in zf.namelist():
            root = ET.fromstring(zf.read("Contents/content.hpf"))
            for item in (node for node in root.iter() if _local(node.tag) == "item"):
                if item.attrib.get("id") and item.attrib.get("href"):
                    manifest[item.attrib["id"]] = "Contents/" + item.attrib["href"].lstrip("/")
        blocks = []
        counter = 0
        for section in sorted(name for name in zf.namelist() if name.startswith("Contents/section") and name.endswith(".xml")):
            root = ET.fromstring(zf.read(section))
            for para in (node for node in root.iter() if _local(node.tag) == "p"):
                text = " ".join((node.text or "").strip() for node in para.iter() if _local(node.tag) == "t" and (node.text or "").strip())
                images = []
                for img in (node for node in para.iter() if _local(node.tag) == "img"):
                    binary_id = next((value for key, value in img.attrib.items() if _local(key) == "binaryItemIDRef"), None)
                    target = manifest.get(binary_id or "")
                    if target and target in zf.namelist():
                        counter += 1
                        placeholder = f"[[HWPX_FIGURE_{counter}]]"
                        images.append({"data": zf.read(target), "name": target, "placeholder": placeholder, "location": text or section, "matchMethod": "hwpx_binary_item_reference", "matchConfidence": 1.0})
                        text = f"{text}\n{placeholder}".strip()
                if text or images: blocks.append({"text": text, "images": images})
    generated = "\n\n".join(block["text"] for block in blocks)
    return _finalize(data=data, filename="document.hwpx", format_name="hwpx", base_text=generated or base_text, blocks=blocks, assets_root=assets_root, image_analyzer=image_analyzer)


def extract_pptx_figures(data: bytes, base_text: str, assets_root: Path, image_analyzer=None) -> dict:
    presentation = Presentation(io.BytesIO(data))
    blocks = []
    counter = 0
    for slide_no, slide in enumerate(presentation.slides, 1):
        for shape_no, shape in enumerate(slide.shapes, 1):
            text = str(getattr(shape, "text", "") or "").strip()
            images = []
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counter += 1
                placeholder = f"[[PPTX_FIGURE_{counter}]]"
                images.append({"data": shape.image.blob, "name": f"slide{slide_no}-{shape.image.filename}", "placeholder": placeholder, "location": f"슬라이드 {slide_no}, 객체 {shape_no}: {text}", "matchMethod": "pptx_slide_shape_order", "matchConfidence": 1.0})
                text = f"{text}\n{placeholder}".strip()
            if text or images: blocks.append({"text": f"[슬라이드 {slide_no}] {text}".strip(), "images": images})
    generated = "\n\n".join(block["text"] for block in blocks)
    return _finalize(data=data, filename="document.pptx", format_name="pptx", base_text=generated or base_text, blocks=blocks, assets_root=assets_root, image_analyzer=image_analyzer)


def extract_xlsx_figures(data: bytes, base_text: str, assets_root: Path, image_analyzer=None) -> dict:
    workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    blocks = []
    counter = 0
    for sheet in workbook.worksheets:
        for image in getattr(sheet, "_images", []):
            counter += 1
            anchor = getattr(image, "anchor", None)
            row = int(getattr(getattr(anchor, "_from", None), "row", 0)) + 1
            column = int(getattr(getattr(anchor, "_from", None), "col", 0)) + 1
            nearby = []
            for values in sheet.iter_rows(min_row=max(1, row - 1), max_row=row + 1, values_only=True):
                line = " | ".join(str(value).strip() for value in values if value not in (None, ""))
                if line: nearby.append(line)
            placeholder = f"[[XLSX_FIGURE_{counter}]]"
            location = f"Sheet {sheet.title}, {sheet.cell(row=row, column=column).coordinate} 근처: {' / '.join(nearby)}"
            blocks.append({"text": f"[Sheet: {sheet.title}]\n{location}\n{placeholder}", "images": [{"data": image._data(), "name": f"{sheet.title}-{counter}.{getattr(image, 'format', 'png')}", "placeholder": placeholder, "location": location, "matchMethod": "xlsx_drawing_anchor", "matchConfidence": 1.0}]})
    generated = base_text + ("\n\n" + "\n\n".join(block["text"] for block in blocks) if blocks else "")
    return _finalize(data=data, filename="document.xlsx", format_name="xlsx", base_text=generated, blocks=blocks, assets_root=assets_root, image_analyzer=image_analyzer)


def extract_pdf_figures(data: bytes, base_text: str, assets_root: Path, image_analyzer=None) -> dict:
    reader = PdfReader(io.BytesIO(data))
    blocks = []
    counter = 0
    for page_no, page in enumerate(reader.pages, 1):
        page_text = str(page.extract_text() or "").strip()
        images = []
        for image in page.images:
            counter += 1
            placeholder = f"[[PDF_FIGURE_{counter}]]"
            images.append({"data": image.data, "name": image.name, "placeholder": placeholder, "location": f"PDF {page_no}페이지: {_clip(page_text, 500)}", "matchMethod": "pdf_page_image_object", "matchConfidence": 0.8})
        if images:
            markers = "\n".join(raw["placeholder"] for raw in images)
            blocks.append({"text": f"[PDF {page_no}페이지]\n{page_text}\n{markers}", "images": images})
    generated = base_text + ("\n\n## PDF 페이지별 그림\n\n" + "\n\n".join(block["text"] for block in blocks) if blocks else "")
    return _finalize(data=data, filename="document.pdf", format_name="pdf", base_text=generated, blocks=blocks, assets_root=assets_root, image_analyzer=image_analyzer)
